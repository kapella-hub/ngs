"""Document ingestion and text extraction."""
import logging
import shutil
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Iterable, Set, Tuple, Optional
from datetime import datetime
from urllib.parse import urljoin, urlparse

# Document processing imports
import pypdf
from docx import Document
from pptx import Presentation
import pandas as pd
from bs4 import BeautifulSoup
from PIL import Image
import pytesseract
import requests

from app.chunking import chunk_text, chunk_text_smart, chunk_text_with_paragraphs
from app.vectorstore import add_chunks
from app.config import settings

logger = logging.getLogger(__name__)


def chunk_text_by_strategy(text: str) -> List[str]:
    """
    Chunk text using the configured chunking strategy.
    
    Args:
        text: Text to chunk
        
    Returns:
        List of text chunks
        
    Note:
        For paragraph strategy, extracts only the text from chunks with metadata.
    """
    strategy = settings.chunking_strategy
    
    logger.debug(f"Using chunking strategy: {strategy}")
    
    if strategy == "simple":
        return chunk_text(text)
    elif strategy == "paragraph":
        # Paragraph strategy returns dicts with metadata, extract text only
        chunks_with_meta = chunk_text_with_paragraphs(text)
        return [chunk["text"] for chunk in chunks_with_meta]
    else:  # "smart" is default
        return chunk_text_smart(text)


def extract_text_from_txt(file_path: Path) -> str:
    """Extract text from .txt or .md file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # Try with different encoding
        with open(file_path, 'r', encoding='latin-1') as f:
            return f.read()


def extract_text_from_pdf(file_path: Path) -> str:
    """Extract text from PDF file."""
    try:
        text_parts = []
        with open(file_path, 'rb') as f:
            pdf_reader = pypdf.PdfReader(f)
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"Error extracting text from PDF {file_path}: {e}")
        raise


def extract_text_from_docx(file_path: Path) -> str:
    """Extract text from Word document."""
    try:
        doc = Document(file_path)
        text_parts = []
        
        # Extract from paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)
        
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {file_path}: {e}")
        raise


def extract_text_from_pptx(file_path: Path) -> str:
    """Extract text from PowerPoint presentation."""
    try:
        prs = Presentation(file_path)
        text_parts = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {file_path}: {e}")
        raise


def extract_text_from_html(file_path: Path) -> str:
    """Extract text from HTML file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'lxml')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Get text
        text = soup.get_text(separator="\n")
        
        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        return text
    except Exception as e:
        logger.error(f"Error extracting text from HTML {file_path}: {e}")
        raise


def extract_text_from_csv(file_path: Path) -> str:
    """Extract text from CSV file."""
    try:
        df = pd.read_csv(file_path)
        
        # Convert dataframe to text representation
        text_parts = []
        
        # Add column headers
        text_parts.append(" | ".join(df.columns))
        text_parts.append("-" * 50)
        
        # Add rows
        for _, row in df.iterrows():
            row_text = " | ".join(str(val) for val in row.values)
            text_parts.append(row_text)
        
        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"Error extracting text from CSV {file_path}: {e}")
        raise


def extract_text_from_excel(file_path: Path) -> str:
    """Extract text from Excel file (.xls, .xlsx)."""
    try:
        # Read all sheets
        excel_file = pd.ExcelFile(file_path)
        text_parts = []
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            text_parts.append(f"\n=== Sheet: {sheet_name} ===\n")
            
            # Add column headers
            text_parts.append(" | ".join(df.columns.astype(str)))
            text_parts.append("-" * 50)
            
            # Add rows
            for _, row in df.iterrows():
                row_text = " | ".join(str(val) for val in row.values)
                text_parts.append(row_text)
        
        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"Error extracting text from Excel {file_path}: {e}")
        raise


def save_image_file(image_path: Path, doc_id: str, image_index: int = 0) -> str:
    """
    Save an image file to the configured image directory with a unique name.
    
    Args:
        image_path: Path to the original image file
        doc_id: Logical document ID (used for unique naming)
        image_index: Index for multiple images from same document (default: 0)
        
    Returns:
        Relative path to the saved image (relative to data directory)
    """
    try:
        # Get image extension
        extension = image_path.suffix.lower()
        
        # Create unique filename: doc_id_index.ext
        # Replace any path separators in doc_id to avoid directory issues
        safe_doc_id = doc_id.replace("/", "_").replace("\\", "_").replace(":", "_")
        image_filename = f"{safe_doc_id}_{image_index}{extension}"
        
        # Full path to save location
        save_path = Path(settings.image_dir) / image_filename
        
        # Copy image to storage directory
        shutil.copy2(image_path, save_path)
        
        logger.debug(f"Saved image: {image_filename}")
        
        # Return relative path (for storage in metadata)
        return str(save_path.relative_to(settings.data_dir))
        
    except Exception as e:
        logger.error(f"Error saving image {image_path}: {e}")
        raise


def extract_text_from_image(file_path: Path) -> str:
    """
    Extract text from image using OCR (pytesseract).
    
    Note: Requires Tesseract to be installed on the system:
    - Windows: Download installer from https://github.com/UB-Mannheim/tesseract/wiki
    - Linux: sudo apt-get install tesseract-ocr
    - macOS: brew install tesseract
    """
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        return text
    except pytesseract.TesseractNotFoundError:
        logger.error(
            "Tesseract is not installed. Please install Tesseract OCR:\n"
            "  Windows: https://github.com/UB-Mannheim/tesseract/wiki\n"
            "  Linux: sudo apt-get install tesseract-ocr\n"
            "  macOS: brew install tesseract"
        )
        raise
    except Exception as e:
        logger.error(f"Error extracting text from image {file_path}: {e}")
        raise


def extract_links_from_html(soup: BeautifulSoup, base_url: str) -> List[str]:
    """
    Extract all links from HTML content.
    
    Args:
        soup: BeautifulSoup object with parsed HTML
        base_url: Base URL to resolve relative links
        
    Returns:
        List of absolute URLs found in the page
    """
    links = []
    
    # Find all <a> tags with href attribute
    for link in soup.find_all('a', href=True):
        href = link['href'].strip()
        
        # Skip empty links, anchors, javascript, mailto, etc.
        if not href or href.startswith('#') or href.startswith('javascript:') or \
           href.startswith('mailto:') or href.startswith('tel:'):
            continue
        
        # Convert relative URLs to absolute
        absolute_url = urljoin(base_url, href)
        
        # Only include http/https URLs
        parsed = urlparse(absolute_url)
        if parsed.scheme in ('http', 'https'):
            links.append(absolute_url)
    
    return links


def extract_text_from_url(url: str, timeout: int = 30) -> Tuple[str, List[str]]:
    """
    Extract text and links from a web page URL or PDF document.
    
    Args:
        url: URL to fetch
        timeout: Request timeout in seconds (default: 30)
        
    Returns:
        Tuple of (extracted text content, list of links found on the page)
        For PDFs, the links list will be empty.
        
    Raises:
        requests.RequestException: If fetching the URL fails
        ValueError: If the content type is not supported or text extraction fails
    """
    try:
        logger.info(f"Fetching content from URL: {url}")
        
        # Fetch the URL
        headers = {
            'User-Agent': 'Mozilla/5.0 (compatible; RAG-System/1.0)'
        }
        response = requests.get(url, headers=headers, timeout=timeout)
        response.raise_for_status()
        
        # Check content type
        content_type = response.headers.get('content-type', '').lower()
        
        # Handle PDF content
        if 'application/pdf' in content_type or url.lower().endswith('.pdf'):
            logger.info(f"Detected PDF content from {url}, downloading and extracting text...")
            
            # Create a temporary file to store the PDF
            temp_file = None
            try:
                with tempfile.NamedTemporaryFile(mode='wb', suffix='.pdf', delete=False) as temp_file:
                    temp_file.write(response.content)
                    temp_file_path = temp_file.name
                
                # Extract text from the PDF using existing function
                text = extract_text_from_pdf(Path(temp_file_path))
                
                if not text or not text.strip():
                    raise ValueError("No text content extracted from PDF")
                
                logger.info(f"Successfully extracted {len(text)} characters from PDF: {url}")
                # PDFs don't have HTML links, return empty list
                return text, []
                
            finally:
                # Clean up temporary file
                if temp_file and Path(temp_file_path).exists():
                    try:
                        Path(temp_file_path).unlink()
                        logger.debug(f"Cleaned up temporary PDF file: {temp_file_path}")
                    except Exception as cleanup_error:
                        logger.warning(f"Failed to clean up temporary file {temp_file_path}: {cleanup_error}")
        
        # Handle HTML/text content
        elif 'html' in content_type or 'text' in content_type:
            # Parse HTML and extract text
            soup = BeautifulSoup(response.content, 'html.parser')
            
            # Extract links before removing script/style
            links = extract_links_from_html(soup, url)
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            if not text or not text.strip():
                raise ValueError("No text content extracted from URL")
            
            logger.info(f"Successfully extracted {len(text)} characters and {len(links)} links from {url}")
            return text, links
        
        else:
            # Unsupported content type
            raise ValueError(
                f"Unsupported content type: {content_type}. "
                f"Supported types: HTML, text, and PDF (application/pdf)."
            )
        
    except requests.Timeout:
        logger.error(f"Timeout while fetching URL: {url}")
        raise
    except requests.RequestException as e:
        logger.error(f"Error fetching URL {url}: {e}")
        raise
    except Exception as e:
        logger.error(f"Error extracting text from URL {url}: {e}")
        raise


def crawl_url_with_links(
    start_url: str,
    max_depth: int = 1,
    same_domain_only: bool = True,
    timeout: int = 30,
    cancellation_event: Optional[Any] = None
) -> Dict[str, Any]:
    """
    Crawl a URL and optionally follow links found on the page.
    
    Args:
        start_url: Starting URL to crawl
        max_depth: Maximum depth to crawl (1 = only start URL, 2 = start + linked pages, etc.)
        same_domain_only: If True, only follow links within the same domain
        timeout: Request timeout in seconds for each URL
        cancellation_event: Optional threading.Event to check for cancellation requests
        
    Returns:
        Dictionary with crawling results:
            - urls_processed: List of URLs that were successfully processed
            - urls_failed: List of URLs that failed to process
            - total_text_length: Total characters extracted from all pages
            - all_texts: Dictionary mapping URL to extracted text
            - cancelled: True if operation was cancelled
    """
    visited: Set[str] = set()
    to_visit: List[Tuple[str, int]] = [(start_url, 0)]  # (url, depth)
    
    urls_processed = []
    urls_failed = []
    all_texts = {}
    cancelled = False
    
    # Parse start URL domain
    start_domain = urlparse(start_url).netloc
    
    while to_visit:
        # Check for cancellation at the start of each iteration
        if cancellation_event and cancellation_event.is_set():
            logger.info(f"Crawling cancelled by user after processing {len(urls_processed)} URLs")
            cancelled = True
            break
        
        current_url, current_depth = to_visit.pop(0)
        
        # Skip if already visited
        if current_url in visited:
            continue
        
        # Skip if depth exceeded
        if current_depth >= max_depth:
            continue
        
        # Skip if domain restriction is enabled and domain doesn't match
        if same_domain_only:
            current_domain = urlparse(current_url).netloc
            if current_domain != start_domain:
                logger.debug(f"Skipping {current_url} - different domain")
                continue
        
        visited.add(current_url)
        
        try:
            # Extract text and links from the URL
            text, links = extract_text_from_url(current_url, timeout=timeout)
            
            if text and text.strip():
                all_texts[current_url] = text
                urls_processed.append(current_url)
                logger.info(f"Successfully processed {current_url} (depth {current_depth})")
                
                # Add discovered links to the queue if we haven't reached max depth
                if current_depth + 1 < max_depth:
                    for link in links:
                        if link not in visited:
                            to_visit.append((link, current_depth + 1))
                            logger.debug(f"Queued {link} for crawling at depth {current_depth + 1}")
            else:
                logger.warning(f"No content extracted from {current_url}")
                urls_failed.append(current_url)
                
        except Exception as e:
            logger.error(f"Failed to process {current_url}: {e}")
            urls_failed.append(current_url)
    
    total_text_length = sum(len(text) for text in all_texts.values())
    
    if cancelled:
        logger.info(
            f"Crawling cancelled: {len(urls_processed)} URLs processed before cancellation, "
            f"{len(urls_failed)} failed, {total_text_length} total characters"
        )
    else:
        logger.info(
            f"Crawling complete: {len(urls_processed)} URLs processed, "
            f"{len(urls_failed)} failed, {total_text_length} total characters"
        )
    
    return {
        "urls_processed": urls_processed,
        "urls_failed": urls_failed,
        "total_text_length": total_text_length,
        "all_texts": all_texts,
        "cancelled": cancelled
    }


def extract_text(file_path: Path) -> str:
    """
    Extract text from a file based on its extension.
    
    Args:
        file_path: Path to the file
        
    Returns:
        Extracted text content
        
    Raises:
        ValueError: If file type is not supported
    """
    suffix = file_path.suffix.lower()
    
    extractors = {
        '.txt': extract_text_from_txt,
        '.md': extract_text_from_txt,
        '.pdf': extract_text_from_pdf,
        '.docx': extract_text_from_docx,
        '.pptx': extract_text_from_pptx,
        '.html': extract_text_from_html,
        '.htm': extract_text_from_html,
        '.csv': extract_text_from_csv,
        '.xls': extract_text_from_excel,
        '.xlsx': extract_text_from_excel,
        '.png': extract_text_from_image,
        '.jpg': extract_text_from_image,
        '.jpeg': extract_text_from_image,
        '.tif': extract_text_from_image,
        '.tiff': extract_text_from_image,
    }
    
    extractor = extractors.get(suffix)
    if not extractor:
        raise ValueError(f"Unsupported file type: {suffix}")
    
    logger.info(f"Extracting text from {file_path.name} ({suffix})")
    return extractor(file_path)


def ingest_files(
    file_paths: Iterable[Path],
    logical_doc_id_prefix: str
) -> Dict[str, Any]:
    """
    Ingest multiple files into the vector store.
    
    Args:
        file_paths: Iterable of file paths to ingest
        logical_doc_id_prefix: Prefix for document IDs (e.g., batch_id)
        
    Returns:
        Dictionary with ingestion results:
            - success_count: Number of successfully ingested files
            - failed_count: Number of failed files
            - total_chunks: Total number of chunks created
            - failed_files: List of failed file names with error messages
    """
    success_count = 0
    failed_count = 0
    total_chunks = 0
    failed_files = []
    
    all_texts = []
    all_doc_ids = []
    all_sources = []
    all_image_paths = []  # Track image paths for each chunk
    
    # Image file extensions
    image_extensions = {'.png', '.jpg', '.jpeg', '.tif', '.tiff'}
    
    for idx, file_path in enumerate(file_paths):
        try:
            # Extract text
            text = extract_text(file_path)
            
            if not text or not text.strip():
                logger.warning(f"No text extracted from {file_path.name}")
                failed_count += 1
                failed_files.append({
                    "file": file_path.name,
                    "error": "No text content extracted"
                })
                continue
            
            # Chunk text using configured strategy
            chunks = chunk_text_by_strategy(text)
            
            if not chunks:
                logger.warning(f"No chunks created from {file_path.name}")
                failed_count += 1
                failed_files.append({
                    "file": file_path.name,
                    "error": "No chunks created"
                })
                continue
            
            # Prepare metadata
            doc_id = f"{logical_doc_id_prefix}:{idx}"
            source = file_path.name
            
            # Check if this is an image file - if so, save it
            image_path = None
            if file_path.suffix.lower() in image_extensions:
                try:
                    image_path = save_image_file(file_path, doc_id, image_index=0)
                    logger.info(f"Saved image for {file_path.name}: {image_path}")
                except Exception as e:
                    logger.warning(f"Failed to save image {file_path.name}: {e}")
                    # Continue processing even if image save fails
            
            # Collect chunks
            all_texts.extend(chunks)
            all_doc_ids.extend([doc_id] * len(chunks))
            all_sources.extend([source] * len(chunks))
            # Associate image path with all chunks from this image (or None for non-images)
            all_image_paths.extend([image_path] * len(chunks))
            
            total_chunks += len(chunks)
            success_count += 1
            
            logger.info(f"Successfully processed {file_path.name}: {len(chunks)} chunks")
            
        except Exception as e:
            logger.error(f"Failed to process {file_path.name}: {e}")
            failed_count += 1
            failed_files.append({
                "file": file_path.name,
                "error": str(e)
            })
    
    # Add all chunks to vector store
    if all_texts:
        try:
            add_chunks(all_texts, all_doc_ids, all_sources, image_paths=all_image_paths)
            logger.info(f"Added {len(all_texts)} chunks to vector store")
        except Exception as e:
            logger.error(f"Failed to add chunks to vector store: {e}")
            raise
    
    return {
        "success_count": success_count,
        "failed_count": failed_count,
        "total_chunks": total_chunks,
        "failed_files": failed_files
    }


def ingest_url(
    url: str,
    logical_doc_id: str,
    follow_links: bool = False,
    max_depth: int = 1,
    same_domain_only: bool = True,
    job_id: Optional[str] = None
) -> Dict[str, Any]:
    """
    Ingest content from a URL into the vector store.
    Optionally crawl and ingest linked pages as well.
    
    Args:
        url: URL to fetch and ingest
        logical_doc_id: Unique document ID for this URL
        follow_links: If True, crawl and ingest linked pages
        max_depth: Maximum crawl depth (1 = only start URL, 2 = start + linked pages)
        same_domain_only: If True, only follow links within the same domain
        job_id: Optional job ID for cancellation tracking
        
    Returns:
        Dictionary with ingestion results:
            - success: Boolean indicating success
            - total_chunks: Number of chunks created
            - urls_processed: Number of URLs successfully processed
            - urls_failed: Number of URLs that failed
            - cancelled: True if operation was cancelled
            - error: Error message if failed (optional)
            - source_url: The source URL
            - last_fetched: ISO timestamp of when content was fetched
    """
    # Get cancellation event if job_id is provided
    cancellation_event = None
    if job_id:
        from app.cancellation import cancellation_store
        cancellation_event = cancellation_store.create_job(job_id)
    
    try:
        # Get current timestamp
        last_fetched = datetime.utcnow().isoformat()
        
        if follow_links and max_depth > 1:
            # Crawl multiple URLs
            logger.info(f"Crawling {url} with max_depth={max_depth}, same_domain_only={same_domain_only}")
            crawl_result = crawl_url_with_links(url, max_depth, same_domain_only, cancellation_event=cancellation_event)
            
            # Check if operation was cancelled
            if crawl_result.get("cancelled", False):
                logger.info(f"URL ingestion cancelled for {url}")
                return {
                    "success": False,
                    "total_chunks": 0,
                    "urls_processed": len(crawl_result["urls_processed"]),
                    "urls_failed": len(crawl_result["urls_failed"]),
                    "cancelled": True,
                    "error": "Operation cancelled by user",
                    "source_url": url
                }
            
            if not crawl_result["urls_processed"]:
                return {
                    "success": False,
                    "total_chunks": 0,
                    "urls_processed": 0,
                    "urls_failed": len(crawl_result["urls_failed"]),
                    "cancelled": False,
                    "error": "No URLs were successfully processed",
                    "source_url": url
                }
            
            # Process all crawled texts
            all_chunks = []
            all_doc_ids = []
            all_sources = []
            
            for crawled_url, text in crawl_result["all_texts"].items():
                # Chunk text using configured strategy
                chunks = chunk_text_by_strategy(text)
                
                if chunks:
                    all_chunks.extend(chunks)
                    # Use unique doc ID for each URL
                    url_doc_id = f"{logical_doc_id}:{crawled_url}"
                    all_doc_ids.extend([url_doc_id] * len(chunks))
                    all_sources.extend([crawled_url] * len(chunks))
            
            if not all_chunks:
                return {
                    "success": False,
                    "total_chunks": 0,
                    "urls_processed": len(crawl_result["urls_processed"]),
                    "urls_failed": len(crawl_result["urls_failed"]),
                    "cancelled": False,
                    "error": "No chunks created from crawled content",
                    "source_url": url
                }
            
            # Add all chunks to vector store
            add_chunks(all_chunks, all_doc_ids, all_sources, upload_date=last_fetched)
            
            logger.info(
                f"Successfully ingested {len(crawl_result['urls_processed'])} URLs "
                f"from {url}: {len(all_chunks)} total chunks"
            )
            
            return {
                "success": True,
                "total_chunks": len(all_chunks),
                "urls_processed": len(crawl_result["urls_processed"]),
                "urls_failed": len(crawl_result["urls_failed"]),
                "cancelled": False,
                "source_url": url,
                "last_fetched": last_fetched
            }
        else:
            # Single URL ingestion (original behavior)
            text, links = extract_text_from_url(url)
            
            if not text or not text.strip():
                return {
                    "success": False,
                    "total_chunks": 0,
                    "urls_processed": 0,
                    "urls_failed": 1,
                    "cancelled": False,
                    "error": "No text content extracted from URL",
                    "source_url": url
                }
            
            # Chunk text using configured strategy
            chunks = chunk_text_by_strategy(text)
            
            if not chunks:
                return {
                    "success": False,
                    "total_chunks": 0,
                    "urls_processed": 0,
                    "urls_failed": 1,
                    "cancelled": False,
                    "error": "No chunks created from URL content",
                    "source_url": url
                }
            
            # Prepare metadata
            doc_ids = [logical_doc_id] * len(chunks)
            sources = [url] * len(chunks)
            
            # Add chunks to vector store with URL metadata
            add_chunks(chunks, doc_ids, sources, upload_date=last_fetched)
            
            logger.info(f"Successfully ingested URL {url}: {len(chunks)} chunks")
            
            return {
                "success": True,
                "total_chunks": len(chunks),
                "urls_processed": 1,
                "urls_failed": 0,
                "cancelled": False,
                "source_url": url,
                "last_fetched": last_fetched
            }
        
    except Exception as e:
        logger.error(f"Failed to ingest URL {url}: {e}")
        return {
            "success": False,
            "total_chunks": 0,
            "urls_processed": 0,
            "urls_failed": 1,
            "cancelled": False,
            "error": str(e),
            "source_url": url
        }
    finally:
        # Cleanup job from cancellation store
        if job_id:
            from app.cancellation import cancellation_store
            cancellation_store.cleanup_job(job_id)


def refresh_url_content(url: str, logical_doc_id: str) -> Dict[str, Any]:
    """
    Refresh content from a URL by re-fetching and updating the vector store.
    
    This function will:
    1. Delete existing chunks for this URL
    2. Fetch fresh content from the URL
    3. Re-chunk and re-index the new content
    
    Args:
        url: URL to refresh
        logical_doc_id: Document ID for this URL
        
    Returns:
        Dictionary with refresh results:
            - success: Boolean indicating success
            - total_chunks: Number of chunks created
            - error: Error message if failed (optional)
            - source_url: The source URL
            - last_fetched: ISO timestamp of when content was fetched
    """
    from app.vectorstore import delete_file_by_source
    
    try:
        # Delete existing chunks for this URL
        delete_result = delete_file_by_source(url)
        logger.info(f"Deleted {delete_result.get('deleted_count', 0)} existing chunks for {url}")
        
        # Re-ingest the URL
        result = ingest_url(url, logical_doc_id)
        
        if result["success"]:
            logger.info(f"Successfully refreshed URL {url}")
        
        return result
        
    except Exception as e:
        logger.error(f"Failed to refresh URL {url}: {e}")
        return {
            "success": False,
            "total_chunks": 0,
            "error": str(e),
            "source_url": url
        }
